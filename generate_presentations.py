"""Generate OCP Java 21 Chapter 1 & 2 presentations as .pptx for Google Slides import."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy body
BG_TITLE  = RGBColor(0x13, 0x13, 0x2A)   # even darker for title slides
ACCENT    = RGBColor(0xF3, 0x8B, 0x1A)   # orange accent
ACCENT2   = RGBColor(0x5B, 0xD8, 0xB8)   # teal for code
RED_TRAP  = RGBColor(0xFF, 0x5C, 0x5C)   # exam trap red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY= RGBColor(0xCC, 0xCC, 0xCC)
CODE_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # near-black for code bg
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = font_name
    return txBox


def add_label(slide, text, left, top, width, height, **kwargs):
    return add_textbox(slide, text, left, top, width, height, **kwargs)


def code_box(slide, code_text, left, top, width, height, font_size=13):
    """Dark-background monospaced code box."""
    add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.07),
        width - Inches(0.2), height - Inches(0.14)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.name  = "Courier New"
        run.font.color.rgb = ACCENT2
    return txBox


def trap_box(slide, text, left, top, width, height, font_size=15):
    """Red-bordered warning/trap box."""
    add_rect(slide, left, top, width, height,
             fill_color=RGBColor(0x3A, 0x10, 0x10), line_color=RED_TRAP)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.1),
        width - Inches(0.24), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "EXAM TRAP:  " + text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RED_TRAP
    run.font.name = "Calibri"
    return txBox


def section_title_slide(prs, chapter, chapter_title, section):
    """Full-screen section divider."""
    slide = blank_slide(prs)
    fill_bg(slide, BG_TITLE)
    add_rect(slide, 0, Inches(3.0), SLIDE_W, Inches(0.06), fill_color=ACCENT)
    add_textbox(slide, chapter, Inches(0.5), Inches(1.0),
                Inches(12), Inches(0.7),
                font_size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_textbox(slide, chapter_title, Inches(0.5), Inches(1.9),
                Inches(12), Inches(1.0),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, section, Inches(0.5), Inches(3.3),
                Inches(12), Inches(0.8),
                font_size=28, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    return slide


def chapter_title_slide(prs, chapter, title, subtitle="OCP Java SE 21 — Exam 1Z0-830"):
    slide = blank_slide(prs)
    fill_bg(slide, BG_TITLE)
    add_rect(slide, 0, Inches(0), SLIDE_W, Inches(0.5), fill_color=ACCENT)
    add_rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.5), fill_color=ACCENT)
    add_textbox(slide, subtitle, Inches(0.5), Inches(0.6),
                Inches(12), Inches(0.5),
                font_size=18, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)
    add_textbox(slide, chapter, Inches(0.5), Inches(1.5),
                Inches(12), Inches(0.8),
                font_size=24, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_textbox(slide, title, Inches(0.5), Inches(2.5),
                Inches(12), Inches(1.5),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Boyarsky & Selikoff  |  Weekly Study Prep Lecture",
                Inches(0.5), Inches(5.5),
                Inches(12), Inches(0.6),
                font_size=16, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title, bullets, chapter=""):
    """Standard slide: title bar + bullet list."""
    slide = blank_slide(prs)
    fill_bg(slide, BG_DARK)
    # top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
    # title
    add_textbox(slide, title, Inches(0.3), Inches(0.15),
                Inches(12.5), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    if chapter:
        add_textbox(slide, chapter, Inches(10.5), Inches(0.15),
                    Inches(2.5), Inches(0.5),
                    font_size=13, color=ACCENT, align=PP_ALIGN.RIGHT)
    # horizontal rule
    add_rect(slide, Inches(0.3), Inches(0.85), Inches(12.5), Inches(0.03),
             fill_color=ACCENT)
    # bullets
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.0),
                                     Inches(12.3), Inches(6.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        indent = "    " * level
        bullet_char = "▪" if level == 0 else "–" if level == 1 else "·"
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{indent}{bullet_char}  {text}"
        run.font.size = Pt(18 - level * 1)
        run.font.color.rgb = WHITE if level == 0 else LIGHT_GREY
        run.font.name = "Calibri"
    return slide


def split_slide(prs, title, bullets, code, chapter="", code_font=13):
    """Left: bullets, Right: code block."""
    slide = blank_slide(prs)
    fill_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
    add_textbox(slide, title, Inches(0.3), Inches(0.15),
                Inches(12.5), Inches(0.7),
                font_size=26, bold=True, color=WHITE)
    if chapter:
        add_textbox(slide, chapter, Inches(10.5), Inches(0.15),
                    Inches(2.5), Inches(0.5),
                    font_size=13, color=ACCENT, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.3), Inches(0.85), Inches(12.5), Inches(0.03),
             fill_color=ACCENT)

    # Left panel bullets
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.0),
                                     Inches(5.6), Inches(6.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        indent = "  " * level
        bullet_char = "▪" if level == 0 else "–"
        run = p.add_run()
        run.text = f"{indent}{bullet_char}  {text}"
        run.font.size = Pt(16 - level)
        run.font.color.rgb = WHITE if level == 0 else LIGHT_GREY
        run.font.name = "Calibri"

    # Right panel code
    code_lines = code.strip('\n').split('\n')
    code_h = max(Inches(1.2), min(Inches(6.1), Pt(code_font + 6) * len(code_lines) + Inches(0.3)))
    code_box(slide, code.strip('\n'),
             Inches(6.2), Inches(1.0),
             Inches(6.8), code_h, font_size=code_font)
    return slide


def trap_slide(prs, title, explanation, code, chapter=""):
    """Red-alert trap slide."""
    slide = blank_slide(prs)
    fill_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=RED_TRAP)
    add_textbox(slide, "⚠  EXAM TRAP — " + title,
                Inches(0.3), Inches(0.15),
                Inches(12.5), Inches(0.7),
                font_size=26, bold=True, color=RED_TRAP)
    if chapter:
        add_textbox(slide, chapter, Inches(10.5), Inches(0.15),
                    Inches(2.5), Inches(0.5),
                    font_size=13, color=ACCENT, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.3), Inches(0.85), Inches(12.5), Inches(0.03),
             fill_color=RED_TRAP)
    add_textbox(slide, explanation, Inches(0.4), Inches(1.0),
                Inches(12.3), Inches(2.2),
                font_size=18, color=WHITE)
    if code:
        code_lines = code.strip('\n').split('\n')
        ch = min(Inches(5.8), Pt(15) * len(code_lines) + Inches(0.4))
        code_box(slide, code.strip('\n'),
                 Inches(0.4), Inches(3.3),
                 Inches(12.3), ch, font_size=14)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  CHAPTER 1 — Building Blocks
# ═══════════════════════════════════════════════════════════════════════════════

def build_ch1(prs):
    CH = "Chapter 1"
    T  = "Building Blocks"

    # ── Title ─────────────────────────────────────────────────────────────────
    chapter_title_slide(prs, CH, T)

    # ── Agenda ────────────────────────────────────────────────────────────────
    content_slide(prs, "Chapter 1 Agenda", [
        (0, "Class & File Structure"),
        (0, "The main() Method"),
        (0, "Packages & Imports"),
        (0, "Constructors & Initialization Order"),
        (0, "Primitive Types & Numeric Literals"),
        (0, "Reference Types vs Primitives"),
        (0, "Wrapper Classes"),
        (0, "Text Blocks"),
        (0, "Identifiers & Declaring Variables"),
        (0, "var — Local Variable Type Inference"),
        (0, "Initializing Variables"),
        (0, "Variable Scope"),
        (0, "Garbage Collection"),
        (0, "Exam Traps & Review Mistakes"),
    ], chapter=CH)

    # ── Section: Class & File Structure ───────────────────────────────────────
    section_title_slide(prs, CH, T, "Class & File Structure")

    split_slide(prs, "Class & File Structure", [
        (0, "One public top-level type per .java file"),
        (0, "Filename must match the public class name"),
        (0, "A file can have multiple classes — at most one public"),
        (0, "Top-level types: class, interface, enum, record, annotation"),
        (1, "Required order inside a file — PIC mnemonic:"),
        (1, "P → package declaration (first, if present)"),
        (1, "I → import statements"),
        (1, "C → class/type declaration"),
        (1, "Fields & methods can appear in any order inside the class"),
    ],
    """// DOES NOT COMPILE
// import before package is wrong
import java.util.*;
package structure;   // package MUST come first

// CORRECT order:
package structure;
import java.util.*;

public class Animal { }
class Helper { }   // also OK — not public""",
    chapter=CH)

    # ── Section: main() ───────────────────────────────────────────────────────
    section_title_slide(prs, CH, T, "The main() Method")

    split_slide(prs, "Writing a main() Method", [
        (0, "Entry point: public static void main(String[] args)"),
        (0, "Must be public and static; return type must be void"),
        (0, "Parameter variations — all valid:"),
        (1, "String[] args"),
        (1, "String args[]"),
        (1, "String... args"),
        (0, "final modifier is optional — still valid"),
        (0, "Parameter name (args, x, friends) doesn't matter — type does"),
        (0, "Command-line args are indexed from args[0]"),
        (0, "static means no object needed to call it"),
    ],
    """// All of these are valid main() signatures:
public static void main(String[] args) { }
public static void main(String args[]) { }
public static void main(String... args) { }
public final static void main(final String[] args) { }

// Running with args:
// java MyClass hello world
// args[0] = "hello", args[1] = "world" """,
    chapter=CH)

    # ── Section: Packages & Imports ───────────────────────────────────────────
    section_title_slide(prs, CH, T, "Packages & Imports")

    split_slide(prs, "Packages & Imports — The Rules", [
        (0, "Wildcard * imports all classes directly in the package"),
        (0, "Does NOT include subpackages, fields, or methods"),
        (0, "Only one wildcard per import; must be at the end"),
        (0, "java.lang is auto-imported (String, System, Object…)"),
        (0, "Classes in the same package need NO import"),
        (0, "Naming conflict: specific import wins over wildcard"),
        (0, "Two explicit imports of same name → compiler error"),
        (0, "Fix: use fully qualified name inline"),
    ],
    """import java.nio.*;           // BAD — can't reach java.nio.file.Files
import java.nio.*.*;         // BAD — only one wildcard
import java.nio.file.Paths.* // BAD — can't import methods

// Conflict resolution:
import java.util.Date;   // specific — WINS
import java.sql.*;        // Date here is ignored

// If both needed:
java.util.Date d1 = new java.util.Date();
java.sql.Date  d2 = new java.sql.Date(0);""",
    chapter=CH)

    content_slide(prs, "Command-Line Compilation & Classpath", [
        (0, "javac options:"),
        (1, "-d <dir>         → output .class files (preserves package structure)"),
        (1, "-cp / -classpath / --class-path <path>   → classpath for compilation"),
        (0, "java options:"),
        (1, "-cp <path> packageb.ClassB   → classpath for running"),
        (0, "Case-sensitive: -D (capital) is WRONG for classpath options"),
        (0, "Example flow:"),
        (1, "javac -d out src/myapp/Main.java"),
        (1, "java -cp out myapp.Main"),
    ], chapter=CH)

    trap_slide(prs, "import does NOT reach subpackages",
        "import java.nio.* does NOT give you access to java.nio.file.Files.\n"
        "Wildcards only import classes directly in that package level.",
        """import java.nio.*;
// CANNOT use Files, Paths, etc. — they're in java.nio.file

import java.nio.file.*;   // This works for Files, Paths, etc.
import java.nio.file.Files;  // Or import specifically""",
        chapter=CH)

    # ── Section: Constructors ─────────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Constructors & Initialization Order")

    trap_slide(prs, "Constructor vs Method — Critical Distinction",
        "A constructor has the SAME NAME as the class and NO return type — not even void.\n"
        "A method with a return type (even void) that matches the class name IS NOT a constructor.\n"
        "It compiles but is never called as a constructor.",
        """public class Chick {
    public Chick() { }           // CONSTRUCTOR — no return type
    public void Chick() { }      // NOT a constructor — regular method!
}

// The void Chick() method compiles, but it's just a method.
// new Chick() calls the CONSTRUCTOR only.""",
        chapter=CH)

    split_slide(prs, "Order of Initialization", [
        (0, "Three phases, always in this order:"),
        (1, "1. Static fields & static initializers (class loading, once)"),
        (1, "2. Instance fields & instance initializers (in order they appear)"),
        (1, "3. Constructor body runs last"),
        (0, "Instance initializer blocks run before the constructor"),
        (0, "Multiple initializer blocks run top-to-bottom"),
        (0, "This ordering is tested heavily on the exam"),
    ],
    """public class Chick {
    private String name = "Fluffy";      // 1st — field init
    { System.out.println("setting field"); }  // 2nd — init block
    public Chick() {
        name = "Tiny";                   // 3rd — constructor
    }
}
// new Chick() output:
// setting field
// (then name = "Tiny")""",
    chapter=CH)

    # ── Section: Primitive Types ──────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Primitive Types")

    content_slide(prs, "The 8 Primitive Types", [
        (0, "boolean  — JVM-dependent size  | true or false   | default: false"),
        (0, "byte     — 8-bit  signed        | -128 to 127     | default: 0"),
        (0, "short    — 16-bit signed        | -32,768 to 32,767 | default: 0"),
        (0, "int      — 32-bit signed        | default integer literal | default: 0"),
        (0, "long     — 64-bit signed        | needs L/l suffix  | default: 0L"),
        (0, "float    — 32-bit IEEE          | NEEDS f/F suffix  | default: 0.0f"),
        (0, "double   — 64-bit IEEE          | default decimal literal | default: 0.0"),
        (0, "char     — 16-bit UNSIGNED      | 0 to 65,535       | default: '\\u0000'"),
        (0, "Key facts:"),
        (1, "boolean's bit size is unspecified — JVM-dependent"),
        (1, "short is signed, char is unsigned — both 16-bit"),
        (1, "float REQUIRES f or F suffix — 3.14 is a double"),
        (1, "Prefer uppercase L for long — lowercase l looks like 1"),
        (1, "String is NOT a primitive — it's a reference type"),
    ], chapter=CH)

    split_slide(prs, "Numeric Literals — Bases & Underscores", [
        (0, "Hexadecimal: 0x or 0X prefix"),
        (0, "Binary: 0b or 0B prefix"),
        (0, "Octal: leading 0"),
        (0, "Underscores for readability — rules:"),
        (1, "Cannot START with underscore"),
        (1, "Cannot END with underscore"),
        (1, "Cannot be adjacent to a decimal point"),
        (1, "Can appear anywhere else (even multiple)"),
    ],
    """long hex    = 0xFF;          // hexadecimal
long binary = 0b1010;        // binary
long octal  = 017;           // octal

int million = 1_000_000;     // OK
double d    = 1_00_0.0_0;    // OK (ugly but compiles)

double bad1 = _1000.0;       // DOES NOT COMPILE — starts with _
double bad2 = 1000.0_;       // DOES NOT COMPILE — ends with _
double bad3 = 1000_.0;       // DOES NOT COMPILE — adjacent to '.'""",
    chapter=CH)

    trap_slide(prs, "float literals MUST have f/F suffix",
        "A bare decimal literal like 3.14 is a double by default.\n"
        "Assigning it to a float WITHOUT the f suffix is a compile error.\n"
        "This is Q23 from your review — a common mistake!",
        """float f = 3.14;    // DOES NOT COMPILE — 3.14 is double
float g = 3.14f;   // OK — explicit float literal
float h = (float) 3.14;  // OK — explicit cast

double d = 3.14;   // fine — decimal defaults to double""",
        chapter=CH)

    # ── Section: Reference Types ──────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Reference Types vs Primitives")

    split_slide(prs, "Reference Types vs Primitives", [
        (0, "Reference types point to an object on the heap"),
        (0, "All references are the same size regardless of object type"),
        (0, "Reference types CAN be null; primitives CANNOT"),
        (0, "Reference types have methods; primitives do not"),
        (0, "Primitive names are lowercase; class names start uppercase"),
        (0, "The OBJECT gets garbage collected, not the reference"),
        (0, "Wrapper classes bridge primitives and objects"),
        (1, "Boolean, Byte, Short, Integer, Long, Float, Double, Character"),
        (1, "Boolean and Character do NOT inherit from Number"),
    ],
    """int value = null;     // DOES NOT COMPILE — primitive
String name = null;   // fine — reference type

// Wrapper parsing:
int x = Integer.parseInt("123");    // returns primitive int
Integer y = Integer.valueOf("123"); // returns Integer object

// Boolean trap:
Boolean.valueOf("kangaroo")  // → false
Boolean.valueOf("true")      // → true (case-insensitive)""",
    chapter=CH)

    # ── Section: Text Blocks ──────────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Text Blocks")

    split_slide(prs, "Text Blocks — Multiline Strings", [
        (0, "Delimiter: triple-quote \"\"\""),
        (0, "Opening \"\"\" MUST be followed by a newline — not on same line as content"),
        (0, "Type is still String — all String methods work"),
        (0, "Essential whitespace: part of the string value"),
        (0, "Incidental whitespace: leading spaces that align code"),
        (1, "Determined by the leftmost non-whitespace character"),
        (1, "Stripped from the output automatically"),
        (0, "Escape sequences:"),
        (1, "\\ at end of line suppresses the newline"),
        (1, "\\s preserves trailing spaces"),
    ],
    """String block = \"\"\"doe\"\"\";    // DOES NOT COMPILE
                              // opening \"\"\" needs newline

String valid = \"\"\"
    doe \\
    deer\"\"\";
// Output: "doe deer"  (backslash suppresses newline)

String poem = \"\"\"
    Roses are red,
    Violets are blue
    \"\"\";
// Incidental spaces stripped; trailing newline included""",
    chapter=CH)

    # ── Section: Identifiers ──────────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Identifiers & var")

    content_slide(prs, "Legal Identifiers — 4 Rules", [
        (0, "1. Must begin with a letter, currency symbol ($, ¥, €), or underscore _"),
        (0, "2. Can include digits but CANNOT start with a digit"),
        (0, "3. Single underscore _ alone is NOT valid (illegal since Java 9)"),
        (0, "4. Cannot use a Java reserved word (true, false, null also off-limits)"),
        (0, "Valid examples:"),
        (1, "okIdentifier    $OK2Identifier    _alsoOK1    __SStillOK$"),
        (0, "Invalid examples:"),
        (1, "3DPointClass   (starts with digit)"),
        (1, "hollywood@vine  (@ not allowed)"),
        (1, "*$coffee        (* not valid start character)"),
        (1, "_               (single underscore)"),
        (1, "public          (reserved word)"),
        (0, "var is NOT a reserved word — can be used as variable/method name"),
        (1, "BUT var IS a reserved type name — cannot be class, interface, or enum name"),
    ], chapter=CH)

    split_slide(prs, "var — Local Variable Type Inference", [
        (0, "var is for LOCAL VARIABLES ONLY. Cannot be used for:"),
        (1, "Instance variables (fields)"),
        (1, "Class/static variables"),
        (1, "Method parameters"),
        (1, "Constructor parameters"),
        (0, "var rules:"),
        (1, "Must be initialized on SAME LINE as declaration"),
        (1, "Type inferred at compile time and CANNOT change"),
        (1, "Cannot be initialized with null (no type to infer)"),
        (1, "Cannot be used in multi-variable declarations"),
    ],
    """public class Zoo {
    var tricky = "Hello";  // DOES NOT COMPILE — instance var

public int add(var a, var b) { // DOES NOT COMPILE — param

var n = null;         // DOES NOT COMPILE — can't infer
var a, b = 3;         // DOES NOT COMPILE — multi-var
var question;         // DOES NOT COMPILE — needs init
question = 1;

var number = 7;       // OK — inferred as int
number = "five";      // DOES NOT COMPILE — already int

var name = "hello";   // String
name = null;          // OK — String can be null""",
    chapter=CH)

    trap_slide(prs, "var in multi-variable declarations — always invalid",
        "var cannot appear in a multi-variable declaration regardless of which variables use var.\n"
        "Q8 and Q18 from your review — you missed this pattern twice!",
        """var a, b = 3;         // DOES NOT COMPILE
int a, var b = 3;     // DOES NOT COMPILE

// var CAN be used as a variable name inside a method:
int var = 3;          // COMPILES — var is not a reserved word
// BUT NOT as a class/interface/enum name:
class var { }         // DOES NOT COMPILE""",
        chapter=CH)

    # ── Section: Initializing Variables ───────────────────────────────────────
    section_title_slide(prs, CH, T, "Initializing Variables")

    split_slide(prs, "Local vs Instance Variable Initialization", [
        (0, "Local variables: NO default value — must initialize before use"),
        (0, "Compiler checks ALL code paths — any uninitialized path = error"),
        (0, "A declared-but-never-used local variable needs no initialization"),
        (0, "Instance & class variables: automatically initialized:"),
        (1, "Objects → null"),
        (1, "Numeric types → 0 / 0L / 0.0f / 0.0"),
        (1, "boolean → false"),
        (1, "char → '\\u0000'"),
        (0, "final local variables cannot be reassigned after initialization"),
    ],
    """int x;
int reply = x + 1;  // DOES NOT COMPILE — x not initialized

// Branch analysis:
int answer;
if (check) { answer = 1; }
else        { answer = 2; }
System.out.println(answer);  // OK — all paths covered

int onlyOne;
if (check) { onlyOne = 1; }
System.out.println(onlyOne);  // DOES NOT COMPILE — else path missing

final int y = 10;
y = 20;   // DOES NOT COMPILE — final""",
    chapter=CH)

    # ── Section: Variable Scope ────────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Variable Scope")

    split_slide(prs, "Variable Scope — By Type", [
        (0, "Local variable — from declaration to END of its block"),
        (0, "Method parameter — for the ENTIRE method"),
        (0, "Instance variable — from declaration until object is eligible for GC"),
        (0, "Class (static) variable — from declaration until PROGRAM ENDS"),
        (0, "Inner blocks can access outer variables — NOT vice versa"),
        (0, "Scope questions often disguised as logic questions"),
        (1, "Always count the closing braces to track scope"),
        (0, "Q12 from review: variable in if block is out of scope after }"),
    ],
    """if (hungry) {
    int bites = 1;
    {
        var tiny = true;    // tiny in scope here
    }
    System.out.println(tiny);  // DOES NOT COMPILE
                               // tiny out of scope after inner }
}
System.out.println(bites);     // DOES NOT COMPILE
                               // bites out of scope after if }""",
    chapter=CH)

    # ── Section: Garbage Collection ───────────────────────────────────────────
    section_title_slide(prs, CH, T, "Garbage Collection")

    split_slide(prs, "Garbage Collection Rules", [
        (0, "All Java objects live on the heap (free store)"),
        (0, "JVM's GC automatically frees memory for unreachable objects"),
        (0, "System.gc() suggests GC — NOT guaranteed to run"),
        (0, "The OBJECT gets GC'd — not the reference"),
        (0, "Object becomes eligible for GC when:"),
        (1, "No reference points to it (all refs reassigned or nulled)"),
        (1, "All references pointing to it have gone out of scope"),
        (0, "Track ALL references — including those passed into methods (Q5!)"),
    ],
    """String one = new String("a");   // one → "a"
String two = new String("b");   // two → "b"
one = two;        // one now → "b"; "a" has NO refs → eligible GC

String three = one;  // three → "b"
one = null;       // "b" still has two and three → NOT eligible

// "b" eligible when both two and three go out of scope
// (end of method)""",
    chapter=CH)

    trap_slide(prs, "Tracking ALL references — Q5 mistake",
        "An object stays alive if ANY reference points to it — including references\n"
        "passed into method calls. Don't only track the named variable.\n"
        "In Q5: brownBear object stayed alive because pandaBear inside roar() still held it.",
        """String brownBear = new BrownBear();   // line 9
roar(brownBear);   // pandaBear inside roar() also points to it!
// ...
brownBear = null;  // line 13

// Object NOT eligible at line 12.
// Only eligible AFTER line 13 AND after roar() returns
// (so pandaBear local ref also goes out of scope)""",
        chapter=CH)

    # ── Section: Review Mistakes ───────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Review Mistakes — Key Reminders")

    content_slide(prs, "Review Mistakes — What to Watch For", [
        (0, "Q3:  main is a METHOD, not a class. bun is a reference, not a class"),
        (0, "Q5:  Track ALL references including those inside called methods"),
        (0, "Q7:  'Which are true' — find ALL correct answers, not just one"),
        (0, "Q8:  var cannot be used in multi-variable declarations — var a, b = 3 fails"),
        (0, "Q11: Same-package classes are automatically visible — NO import needed"),
        (0, "Q12: Variable in if block is out of scope immediately after that block closes"),
        (0, "Q13: Read each option against actual Java rules — not every plausible claim is true"),
        (0, "Q14: var cannot be an instance variable; var var = 3 at field level fails"),
        (0, "Q18: var in multi-variable declaration always invalid — same as Q8"),
        (0, "Q20: A method with a return type (even void) is NEVER a constructor"),
        (0, "Q22: Check all stated constraints before selecting an answer"),
        (0, "Q23: float literals MUST end in f/F — bare decimal is double"),
    ], chapter=CH)

    # ── Summary ───────────────────────────────────────────────────────────────
    slide = blank_slide(prs)
    fill_bg(slide, BG_TITLE)
    add_rect(slide, 0, Inches(0), SLIDE_W, Inches(0.5), fill_color=ACCENT)
    add_textbox(slide, "Chapter 1 — Building Blocks", Inches(0.5), Inches(0.7),
                Inches(12), Inches(0.8), font_size=32, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "Key Takeaways", Inches(0.5), Inches(1.7),
                Inches(12), Inches(0.6), font_size=24, color=ACCENT,
                align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.4), Inches(11), Inches(0.04),
             fill_color=ACCENT)
    key_points = [
        "PIC mnemonic: Package → Import → Class (required file order)",
        "Constructor = same name as class + NO return type",
        "float needs f suffix; long needs L suffix",
        "var = local variables only; must init on same line; no null; no multi-var",
        "Local vars have NO defaults; instance/class vars DO",
        "Wildcard imports don't reach subpackages; java.lang is auto-imported",
        "Object eligible for GC when ALL references gone — track every ref",
        "Single _ is illegal as an identifier (since Java 9)",
    ]
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.6),
                                     Inches(10.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, kp in enumerate(key_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"✓  {kp}"
        run.font.size = Pt(17)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════════
#  CHAPTER 2 — Operators
# ═══════════════════════════════════════════════════════════════════════════════

def build_ch2(prs):
    CH = "Chapter 2"
    T  = "Operators"

    # ── Title ─────────────────────────────────────────────────────────────────
    chapter_title_slide(prs, CH, T)

    # ── Agenda ────────────────────────────────────────────────────────────────
    content_slide(prs, "Chapter 2 Agenda", [
        (0, "Operator Types Overview"),
        (0, "Operator Precedence Table"),
        (0, "Unary Operators (pre-fix vs post-fix)"),
        (0, "Boolean ≠ Numeric — Critical Java Rule"),
        (0, "Compound Assignment + Increment Traps"),
        (0, "Modulo with Negatives"),
        (0, "Numeric Promotion Rules"),
        (0, "Casting — Widening vs Narrowing"),
        (0, "Overflow & Underflow"),
        (0, "Compound Assignment Auto-Cast"),
        (0, "Default Literal Types"),
        (0, "Assignment as an Expression"),
        (0, "instanceof & Pattern Matching"),
        (0, "Logical & Short-Circuit Operators"),
        (0, "Ternary Operator"),
        (0, "Equality Operators"),
        (0, "Exam Traps & Review Mistakes"),
    ], chapter=CH)

    # ── Section: Operator Types ───────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Operator Types & Precedence")

    content_slide(prs, "Three Operator Categories", [
        (0, "Unary — one operand:   !x    ~x    ++x    -x    (int)x"),
        (0, "Binary — two operands:  x + y    x && y    x == y"),
        (0, "Ternary — three operands:  condition ? expr1 : expr2"),
        (0, "Operators evaluated based on PRECEDENCE, not left-to-right"),
        (0, "Same-level operators: left-to-right EXCEPT assignment and ternary"),
        (0, "Assignment and ternary are RIGHT-to-LEFT"),
        (0, "When in doubt: add parentheses and trace step by step"),
    ], chapter=CH)

    content_slide(prs, "Operator Precedence (Highest → Lowest)", [
        (0, "1.  . ()                    method call          Left → Right"),
        (0, "2.  ++ -- (post-fix)                             Left → Right"),
        (0, "3.  ++ -- (pre-fix)  + - (unary)  !  ~  (type)  Right → Left"),
        (0, "4.  *  /  %                 multiply/divide      Left → Right"),
        (0, "5.  +  -                    add/subtract         Left → Right"),
        (0, "6.  <<  >>  >>>             bit shift            Left → Right"),
        (0, "7.  <  >  <=  >=  instanceof relational         Left → Right"),
        (0, "8.  ==  !=                  equality             Left → Right"),
        (0, "9.  &                       bitwise AND          Left → Right"),
        (0, "10. ^                       bitwise XOR          Left → Right"),
        (0, "11. |                       bitwise OR           Left → Right"),
        (0, "12. &&                      conditional AND      Left → Right"),
        (0, "13. ||                      conditional OR       Left → Right"),
        (0, "14. ? :                     ternary              Right → Left"),
        (0, "15. = += -= *= /= %= …      assignment           Right → Left"),
    ], chapter=CH)

    # ── Section: Unary Operators ──────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Unary Operators")

    content_slide(prs, "Unary Operators Reference", [
        (0, "!     → boolean only       — logical inversion"),
        (0, "~     → integral types only — bitwise complement  (~x == -(x+1))"),
        (0, "+ / - → numeric             — unary sign"),
        (0, "++ / -- → numeric           — increment/decrement by 1"),
        (0, "(type) → cast               — explicit type conversion"),
        (0, "CRITICAL: boolean and numeric are COMPLETELY SEPARATE in Java:"),
        (1, "! cannot be applied to numeric types"),
        (1, "~ cannot be applied to boolean"),
        (1, "1 and true have NOTHING to do with each other"),
        (0, "Post-fix has HIGHER precedence than pre-fix"),
    ], chapter=CH)

    trap_slide(prs, "! is boolean-only, ~ is integral-only",
        "Unlike C/C++, Java keeps boolean and numeric types completely separate.\n"
        "You CANNOT apply ! to an integer or ~ to a boolean. This always surprises people.",
        """int x = !5;         // DOES NOT COMPILE — ! is for booleans only
boolean b = ~true;  // DOES NOT COMPILE — ~ is for integrals only

int y = ~5;         // OK — y = -(5+1) = -6
boolean c = !true;  // OK — c = false""",
        chapter=CH)

    split_slide(prs, "Pre-fix vs Post-fix — The Key Rule", [
        (0, "Post-fix (x++): returns ORIGINAL value, THEN increments"),
        (0, "Pre-fix (++x): increments FIRST, THEN returns updated value"),
        (0, "Post-fix has higher precedence than pre-fix"),
        (0, "Trace carefully through any expression with both"),
        (0, "When combined with compound assignment — extra tricky"),
    ],
    """int x = 3;
int y = x++;   // y = 3,  x = 4  (old value returned)
int z = ++x;   // z = 5,  x = 5  (incremented first)

int a = 0;
System.out.println(a++);  // prints 0, a becomes 1
System.out.println(--a);  // a becomes 0, prints 0

// Compound assignment trap:
int k = 0;
k += k++;   // left k=0, k++ returns 0 (k→1), 0+0=0, k=0  → k is 0
k = 0;
k += ++k;   // left k=0, ++k makes k=1 returns 1, 0+1=1  → k is 1""",
    chapter=CH)

    trap_slide(prs, "Compound Assignment captures left side first",
        "Java evaluates operands LEFT TO RIGHT. In compound assignment k += expr,\n"
        "the current value of k is captured BEFORE expr is evaluated.\n"
        "Then the result is written back. This causes very counterintuitive results.",
        """int k = 0;
k += k++;    // Step 1: capture k=0 (left side)
             // Step 2: k++ returns 0, k becomes 1
             // Step 3: 0 + 0 = 0, written to k → k is 0

int k = 0;
k += ++k;    // Step 1: capture k=0 (left side)
             // Step 2: ++k makes k=1, returns 1
             // Step 3: 0 + 1 = 1, written to k → k is 1""",
        chapter=CH)

    split_slide(prs, "Modulo with Negative Numbers", [
        (0, "Result of % takes the SIGN of the LEFT operand (the dividend)"),
        (0, "This is different from mathematical modulo"),
        (0, "Memorize: sign follows the left side"),
        (0, "Examples:"),
        (1, "-13 % 5  → -3  (NOT 2)"),
        (1, " 13 % -5 →  3  (NOT -2)"),
        (1, "-13 % -5 → -3"),
        (1, " 13 %  5 →  3"),
    ],
    """System.out.println(-13 % 5);   // -3, NOT 2
System.out.println(13 % -5);   //  3, NOT -2
System.out.println(-13 % -5);  // -3
System.out.println(13 % 5);    //  3

// Rule: sign of result = sign of LEFT operand (dividend)
// Magnitude = remainder from the division""",
    chapter=CH)

    # ── Section: Numeric Promotion ────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Numeric Promotion & Casting")

    split_slide(prs, "Numeric Promotion Rules", [
        (0, "Rule 1: Smaller type promoted to larger if operands differ in size"),
        (0, "Rule 2: Integral promoted to floating-point if one operand is float/double"),
        (0, "Rule 3: byte, short, char ALWAYS promoted to int with any binary operator"),
        (1, "Even when BOTH operands are the same smaller type!"),
        (0, "Rule 4: Result type = the post-promotion type"),
        (0, "This means byte + byte → int (NOT byte)"),
        (0, "Cannot assign byte + byte result back to byte without cast"),
    ],
    """short x = 14, y = 3;
var z = x * y;    // z is int — both promoted to int first

byte b = 10;
var result = b + b;   // int — NOT byte! Rule 3 applies

byte sum = b + b;     // DOES NOT COMPILE — result is int

byte sum2 = (byte)(b + b);  // OK — explicit cast""",
    chapter=CH)

    split_slide(prs, "Casting — Widening vs Narrowing", [
        (0, "Widening (implicit): safe, compiler handles it automatically"),
        (1, "int → long → float → double"),
        (1, "byte → short → int"),
        (0, "Narrowing (explicit cast): you force it, data MAY be lost"),
        (1, "Floating-point to integral: decimal part TRUNCATED (not rounded)"),
        (1, "Truncation is toward ZERO (so -0.9 → 0, not -1)"),
        (0, "Overflow: value exceeds type range → wraps around silently"),
        (1, "No exception, no warning"),
    ],
    """int i = 5;
double d = i;      // OK — widening, implicit

double pi = 3.14159;
int t = (int) pi;  // t = 3 — decimal DROPPED, NOT rounded

double neg = -0.9;
int n = (int) neg; // n = 0, NOT -1 — toward zero!

byte max = 127;
max++;             // max = -128 — WRAP AROUND, no exception

long a = 100_000_000L;
long b = 100_000_000L;
System.out.println(a * b); // 10000000000000000 — fits in long""",
    chapter=CH)

    trap_slide(prs, "byte+byte = int, and narrowing truncates (not rounds)",
        "Two rules that trip up almost everyone:\n"
        "1. byte + byte → int. You CANNOT assign it back to byte without explicit cast.\n"
        "2. Casting float/double to int TRUNCATES toward zero — no rounding.\n"
        "   (int)(-0.9) = 0 because 0 is closer to zero than -1.",
        """byte a = 10, b = 20;
byte c = a + b;        // DOES NOT COMPILE — result is int
byte c = (byte)(a+b);  // OK

double d = -0.9;
int i = (int) d;  // i = 0, NOT -1
// "toward zero" — the decimal just falls off""",
        chapter=CH)

    # ── Section: Compound Assignment Auto-Cast ────────────────────────────────
    split_slide(prs, "Compound Assignment Operators Auto-Cast", [
        (0, "+=, -=, *=, /=, %= automatically apply narrowing cast back to left-hand type"),
        (0, "Plain assignment does NOT do this"),
        (0, "This is one of the most-tested compound assignment traps"),
        (0, "Equivalence: y *= x  →  y = (int)(y * x)  when y is int"),
    ],
    """long x = 10;
int y = 5;
y *= x;          // COMPILES — implicitly: y = (int)(y * x)
y = y * x;       // DOES NOT COMPILE — explicit cast required
                 //   fix: y = (int)(y * x);

// Same with other compound operators:
byte b = 10;
b += 5;          // COMPILES — auto-cast to byte
b = b + 5;       // DOES NOT COMPILE — b+5 is int""",
    chapter=CH)

    # ── Section: Default Literal Types ────────────────────────────────────────
    split_slide(prs, "Default Literal Types — Common Traps", [
        (0, "Decimal literals (3.14) are DOUBLE by default"),
        (1, "Assigning to float without f suffix → compile error"),
        (0, "Integer literals are INT by default"),
        (1, "Too-large literal for int requires L suffix for long"),
        (0, "These interact with numeric promotion rules"),
    ],
    """float f = 1.1;            // DOES NOT COMPILE — 1.1 is double
float g = 1.1f;           // OK

long z = 1234567890123;   // DOES NOT COMPILE — too big for int
long w = 1234567890123L;  // OK

// Also remember:
int i = 'A';   // OK — char widens to int
// i = 65""",
    chapter=CH)

    # ── Section: Assignment as Expression ────────────────────────────────────
    section_title_slide(prs, CH, T, "Assignment Operators & Equality")

    trap_slide(prs, "Assignment inside if() — always executes!",
        "An assignment returns the value that was assigned.\n"
        "If (boolVar = true) is ASSIGNMENT, not comparison — it ALWAYS executes the body.\n"
        "For non-boolean types, this is a compile error. For boolean, it compiles silently.\n"
        "This is the == vs = confusion — a classic exam gotcha.",
        """boolean healthy = false;
if (healthy = true) {           // COMPILES! Assignment, not ==
    System.out.println("This always prints!");
}

// For non-booleans — compile error:
int x = 5;
if (x = 10) { }   // DOES NOT COMPILE — int is not boolean

// Chaining assignment:
int x = 5;
int y = (x = 10);  // y = 10, x = 10""",
        chapter=CH)

    # ── Section: instanceof ───────────────────────────────────────────────────
    split_slide(prs, "instanceof & Pattern Matching", [
        (0, "a instanceof B — true if a is a B (or subtype of B)"),
        (0, "Pattern matching (Java 16+): auto-casts on match"),
        (1, "if (obj instanceof String s) — s is in scope inside block"),
        (0, "CRITICAL: If it is IMPOSSIBLE for the left side to ever be the right type,"),
        (1, "the COMPILER throws an error — doesn't wait for runtime"),
        (0, "No inheritance relationship between the types → compile error"),
    ],
    """// Pattern matching:
if (obj instanceof String s) {
    System.out.println(s.length());  // s auto-cast to String
}

// Impossible cast — COMPILE ERROR:
Integer i = 5;
if (i instanceof String) { }  // DOES NOT COMPILE
// Integer can NEVER be a String

// Subtype check OK:
Object o = "hello";
if (o instanceof String s) { }  // fine""",
    chapter=CH)

    # ── Section: Logical & Short-Circuit ──────────────────────────────────────
    section_title_slide(prs, CH, T, "Logical & Short-Circuit Operators")

    split_slide(prs, "Logical vs Conditional (Short-Circuit) Operators", [
        (0, "&  (Logical AND)  — ALWAYS evaluates both sides"),
        (0, "|  (Logical OR)   — ALWAYS evaluates both sides"),
        (0, "^  (Logical XOR)  — true only when operands DIFFER"),
        (0, "&& (Conditional AND) — short-circuits: skips right if left is false"),
        (0, "|| (Conditional OR)  — short-circuits: skips right if left is true"),
        (0, "On numeric types: &, |, ^ become BITWISE operators"),
        (0, "Always prefer && over & for null-safe checks"),
        (0, "XOR truth table: A ^ B is true only when A ≠ B"),
    ],
    """String s = null;

// Dangerous — & always evaluates both sides:
if (s != null & s.length() > 0) { }  // NPE!

// Safe — && short-circuits:
if (s != null && s.length() > 0) { } // OK

// XOR:
// true ^ true   = false
// true ^ false  = true
// false ^ true  = true
// false ^ false = false""",
    chapter=CH)

    # ── Section: Ternary ──────────────────────────────────────────────────────
    split_slide(prs, "Ternary Operator", [
        (0, "Syntax: booleanExpression ? exprIfTrue : exprIfFalse"),
        (0, "Only ONE of the two expressions is evaluated"),
        (0, "Both expressions must be compatible with the result type"),
        (0, "Can be used where a value is expected (in assignments, method calls)"),
        (0, "Associativity: right-to-left (same as assignment)"),
        (0, "Nested ternaries are legal but hard to read — watch for them on exam"),
    ],
    """int x = 5;
String result = (x > 3) ? "big" : "small";  // "big"

int y = 1;
int z = (y > 0) ? y++ : y--;  // z = 1, y = 2
                               // only y++ runs

// Nested ternary (legal):
int a = 5;
String s = (a > 10) ? "big" : (a > 3) ? "medium" : "small";
// s = "medium" """,
    chapter=CH)

    # ── Section: Equality ─────────────────────────────────────────────────────
    split_slide(prs, "Equality Operators — == and !=", [
        (0, "For PRIMITIVES: == compares by VALUE"),
        (0, "For REFERENCES: == compares by MEMORY ADDRESS"),
        (1, "Same object? true. Same content? Doesn't matter."),
        (0, "Use .equals() to compare object content"),
        (0, "Numeric promotion applies: int == long → int promoted to long"),
        (0, "null == null → true"),
        (0, "Cannot mix incompatible reference types with =="),
    ],
    """int a = 5, b = 5;
System.out.println(a == b);         // true — same value

String s1 = new String("hi");
String s2 = new String("hi");
System.out.println(s1 == s2);       // false — different objects
System.out.println(s1.equals(s2));  // true  — same content

String s3 = "hello";
String s4 = "hello";
System.out.println(s3 == s4);  // true — string pool, same literal""",
    chapter=CH)

    # ── Section: Review Mistakes ───────────────────────────────────────────────
    section_title_slide(prs, CH, T, "Review Mistakes — Key Reminders")

    content_slide(prs, "Review Mistakes — Chapter 2", [
        (0, "Q3:  Go through EVERY option individually on multi-answer questions"),
        (1, "Answered b, c, f — missed d"),
        (0, "Q8:  Answered e, correct is a — re-trace the arithmetic step by step"),
        (0, "Q10: Answered a, b, d — correct is e"),
        (0, "General — still fuzzy, keep these close:"),
        (1, "Post-fix ++ has HIGHER precedence than pre-fix ++"),
        (1, "byte/short/char ALWAYS promote to int with binary operators — even byte+byte"),
        (1, "Casting float/double to int TRUNCATES — drops decimal, toward zero, NOT round"),
        (1, "Compound assignment (+=, etc.) AUTO-CASTS back to left-hand type — plain = does not"),
        (1, "% result sign follows LEFT operand (dividend)"),
        (1, "if (boolVar = true) compiles and always executes — assignment, not comparison"),
        (1, "& always evaluates both sides; && short-circuits"),
    ], chapter=CH)

    # ── Summary ───────────────────────────────────────────────────────────────
    slide = blank_slide(prs)
    fill_bg(slide, BG_TITLE)
    add_rect(slide, 0, Inches(0), SLIDE_W, Inches(0.5), fill_color=ACCENT)
    add_textbox(slide, "Chapter 2 — Operators", Inches(0.5), Inches(0.7),
                Inches(12), Inches(0.8), font_size=32, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "Key Takeaways", Inches(0.5), Inches(1.7),
                Inches(12), Inches(0.6), font_size=24, color=ACCENT,
                align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, Inches(1.0), Inches(2.4), Inches(11), Inches(0.04),
             fill_color=ACCENT)
    key_points = [
        "Post-fix ++ > Pre-fix ++ in precedence; assignment is right-to-left",
        "! is boolean-only; ~ is integral-only; Java has no truthy/falsy integers",
        "byte/short/char always promote to int with any binary operator",
        "Casting to integral: TRUNCATES toward zero (never rounds)",
        "Compound assignment (+=) auto-casts; plain assignment does not",
        "% sign follows the left operand (dividend)",
        "if (boolVar = true) compiles and always runs — it's assignment!",
        "&& short-circuits; & does not — use && for null-safe checks",
    ]
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.6),
                                     Inches(10.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, kp in enumerate(key_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"✓  {kp}"
        run.font.size = Pt(17)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import os

    # Chapter 1
    prs1 = new_prs()
    build_ch1(prs1)
    out1 = "/home/nashmin/Documents/books/OCP-Prep/chapter1/Chapter1_BuildingBlocks.pptx"
    prs1.save(out1)
    print(f"Saved: {out1}  ({len(prs1.slides)} slides)")

    # Chapter 2
    prs2 = new_prs()
    build_ch2(prs2)
    out2 = "/home/nashmin/Documents/books/OCP-Prep/chapter2/Chapter2_Operators.pptx"
    prs2.save(out2)
    print(f"Saved: {out2}  ({len(prs2.slides)} slides)")
